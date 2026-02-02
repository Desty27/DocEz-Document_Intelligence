import streamlit as st
import os
import tempfile
import shutil
import datetime
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
import plotly.express as px
import faiss
import numpy as np
import json
import base64
from sklearn.metrics import mean_absolute_percentage_error

from indexing import process_documents
from retreival import (
    query_rag_system,
    get_available_models,
    get_processed_files,
    delete_file_from_knowledge_base,
    EMBEDDING_MODELS,
)
from retreival import create_client_for_model, MODELS
from retreival_forecasting import get_forecasting_explanation
from sentence_transformers import SentenceTransformer

import openai

sm = None
try:
    import statsmodels.api as sm

    SARIMAX_AVAILABLE = True
except Exception:
    SARIMAX_AVAILABLE = False

Prophet = None
try:
    from prophet import Prophet

    PROPHET_AVAILABLE = True
except Exception:
    PROPHET_AVAILABLE = False

xgb = None
try:
    import xgboost as xgb

    XGBOOST_AVAILABLE = True
except Exception:
    XGBOOST_AVAILABLE = False

Presentation = None
Inches = None
Pt = None
CategoryChartData = None
XL_CHART_TYPE = None
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False


OPENAI_DEPLOYMENT = "llama3_1_8b_instruct"
OPENAI_BASE_URL = "https://infer.e2enetworks.net/project/p-5915/genai/llama3_1_8b_instruct/v1/"
OPENAI_TOKEN = "eyJhbGciOiJSUzI1NiIsInR5cCIgOiAiSldUIiwia2lkIiA6ICJGSjg2R2NGM2pUYk5MT2NvNE52WmtVQ0lVbWZZQ3FvcXRPUWVNZmJoTmxFIn0.eyJleHAiOjE3ODIxOTIwMjcsImlhdCI6MTc1MDY1NjAyNywianRpIjoiMTdmMWJhZGEtYTYyMS00ZTMwLWJmYTEtZDJhMjYzOTUzMjA4IiwiaXNzIjoiaHR0cDovL2dhdGV3YXkuZTJlbmV0d29ya3MuY29tL2F1dGgvcmVhbG1zL2FwaW1hbiIsImF1ZCI6ImFjY291bnQiLCJzdWIiOiJiMmMzNDZkMC04MmQwLTQxYzItOWVkNS1mNmY4Nzc0MjVlNDkiLCJ0eXAiOiJCZWFyZXIiLCJhenAiOiJhcGltYW51aSIsInNlc3Npb25fc3RhdGUiOiJmNTc3N2QwMi1lYjE4LTRhYjktYmM0NS03ZjVkNmRjZmFjZmUiLCJhY3IiOiIxIiwiYWxsb3dlZC1vcmlnaW5zIjpbIiJdLCJyZWFsbV9hY2Nlc3MiOnsicm9sZXMiOlsib2ZmbGluZV9hY2Nlc3MiLCJ1bWFfYXV0aG9yaXphdGlvbiIsImFwaXVzZXIiLCJkZWZhdWx0LXJvbGVzLWFwaW1hbiJdfSwicmVzb3VyY2VfYWNjZXNzIjp7ImFjY291bnQiOnsicm9sZXMiOlsibWFuYWdlLWFjY291bnQiLCJtYW5hZ2UtYWNjb3VudC1saW5rcyIsInZpZXctcHJvZmlsZSJdfX0sInNjb3BlIjoicHJvZmlsZSBlbWFpbCIsInNpZCI6ImY1Nzc3ZDAyLWViMTgtNGFiOS1iYzQ1LTdmNWQ2ZGNmYWNmZSIsImVtYWlsX3ZlcmlmaWVkIjpmYWxzZSwiaXNfcGFydG5lcl9yb2xlIjpmYWxzZSwibmFtZSI6IkF2aW5hc2ggU2luZ2ggIiwicHJpbWFyeV9lbWFpbCI6InN1cGVyYi5zdWppdEBnbWFpbC5jb20iLCJpc19wcmltYXJ5X2NvbnRhY3QiOmZhbHNlLCJwcmVmZXJyZWRfdXNlcm5hbWUiOiJhdmluYXNoLnNpbmdoLjI3MDUwNEBnbWFpbC5jb20iLCJnaXZlbl9uYW1lIjoiQXZpbmFzaCIsImZhbWlseV9uYW1lIjoiU2luZ2ggIiwiZW1haWwiOiJhdmluYXNoLnNpbmdoLjI3MDUwNEBnbWFpbC5jb20iLCJpc19pbmRpYWFpX3VzZXIiOmZhbHNlfQ.OL2wpdaRRcANjZY9Mx_DQqyoX1_tDAXwqVVGA5rVIGoTV7Bc1hWEy94L0-C3QLTyFJCn1ROd6pbS7rL1qsCcW_XGicfVYLx-PpYFMcAe4GTpxnI_3hn8a_ohjcy8-H6DQ39wXHrVGap8jXsED6V3OurdOw0S_3u8n5bnMw2CIqs"

if OPENAI_TOKEN:
    openai.base_url = OPENAI_BASE_URL  # type: ignore[attr-defined]
    openai.api_key = OPENAI_TOKEN


def create_llm_client():
    """Return an OpenAI-compatible client configured for the E2E endpoint."""
    if not OPENAI_TOKEN:
        raise RuntimeError("API token is not configured.")
    return openai.OpenAI(  # type: ignore[attr-defined]
        api_key=OPENAI_TOKEN,
        base_url=OPENAI_BASE_URL,
    )

# Simple credentials dict
CREDENTIALS = {"gesil": "gesil1234"}

# Page configuration
st.set_page_config(
    page_title="Multimodal GenAI Chatbot System",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for styling
def load_css():
    st.markdown("""
    <style>
    .main {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        min-height: 100vh;
    }
    
    .stApp > header {
        background-color: transparent;
    }
    
    .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
        background: rgba(255, 255, 255, 0.95);
        border-radius: 15px;
        margin: 1rem;
        box-shadow: 0 10px 30px rgba(0, 0, 0, 0.1);
    }
    
    .header-container {
        text-align: center;
        padding: 2rem 0;
        background: linear-gradient(135deg, #FF9736, #DB791D);
        border-radius: 15px;
        margin-bottom: 2rem;
        color: white;
    }
    
    .logo-container {
        display: flex;
        flex-direction: column;
        justify-content: center;
        align-items: center;
        margin-bottom: 1rem;
    }
    
    .company-logo {
        max-height: 80px;
        max-width: 200px;
        margin-bottom: 1rem;
        object-fit: contain;
        filter: drop-shadow(0 2px 4px rgba(0, 0, 0, 0.1));
        transition: transform 0.3s ease;
    }
    
    .company-logo:hover {
        transform: scale(1.05);
    }
    
    .features-container {
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
        gap: 1rem;
        margin-top: 1.5rem;
        padding: 0 2rem;
    }
    
    .feature-item {
        background: rgba(255, 255, 255, 0.15);
        padding: 0.8rem 1.2rem;
        border-radius: 25px;
        font-size: 0.95rem;
        font-weight: 500;
        text-align: center;
        backdrop-filter: blur(10px);
        border: 1px solid rgba(255, 255, 255, 0.2);
        transition: all 0.3s ease;
    }
    
    .feature-item:hover {
        background: rgba(255, 255, 255, 0.25);
        transform: translateY(-2px);
    }
    
    .upload-section {
        background: #f8f9fa;
        padding: 2rem;
        border-radius: 10px;
        border: 2px dashed #3498db;
        margin: 1rem 0;
    }
    
    .query-section {
        background: #f1f3f5;
        padding: 2rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    
    .response-section {
        background: #e8f5e8;
        padding: 2rem;
        border-radius: 10px;
        margin: 1rem 0;
        border-left: 5px solid #27ae60;
    }
    
    .status-success {
        background: #d4edda;
        color: #155724;
        padding: 1rem;
        border-radius: 5px;
        border: 1px solid #c3e6cb;
    }
    
    .status-error {
        background: #f8d7da;
        color: #721c24;
        padding: 1rem;
        border-radius: 5px;
        border: 1px solid #f5c6cb;
    }
    
    .status-info {
        background: #d1ecf1;
        color: #0c5460;
        padding: 1rem;
        border-radius: 5px;
        border: 1px solid #bee5eb;
    }
    
    .metric-card {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.1);
        text-align: center;
        margin: 0.5rem;
    }
    
    .stButton > button {
        background: linear-gradient(135deg, #3498db, #2c3e50);
        color: white;
        border: none;
        border-radius: 25px;
        padding: 0.5rem 2rem;
        font-weight: bold;
        transition: all 0.3s ease;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 5px 15px rgba(0, 0, 0, 0.2);
    }
    
    .sidebar .stButton > button {
        width: 100%;
        margin: 0.5rem 0;
    }
    
    h1, h2, h3 {
        color: #2c3e50;
    }
    
    .title {
        font-size: 3rem;
        font-weight: bold;
        margin-bottom: 0.5rem;
    }
    
    .subtitle {
        font-size: 1.2rem;
        opacity: 0.8;
        margin-bottom: 1rem;
    }
    </style>
    """, unsafe_allow_html=True)

# Initialize session state
def load_logo():
    """Load and encode company logo for display"""
    logo_path = "assets/logo.png"  # You can change this path
    
    # Check if logo file exists
    if os.path.exists(logo_path):
        try:
            with open(logo_path, "rb") as logo_file:
                logo_data = base64.b64encode(logo_file.read()).decode()
            return f"data:image/png;base64,{logo_data}"
        except Exception as e:
            print(f"Error loading logo: {e}")
            return None
    else:
        # If no logo file found, you can use a placeholder or skip
        return None

def upload_logo_interface():
    """Interface for uploading a new company logo"""
    with st.expander("🖼️ Update Company Logo", expanded=False):
        st.write("Upload a new company logo (PNG/JPG, max 2MB)")
        
        uploaded_logo = st.file_uploader(
            "Choose logo file",
            type=['png', 'jpg', 'jpeg'],
            help="Recommended: 200x80 pixels, transparent background PNG"
        )
        
        if uploaded_logo is not None:
            if st.button("Update Logo"):
                try:
                    # Ensure assets directory exists
                    os.makedirs("assets", exist_ok=True)
                    
                    # Save the uploaded logo
                    with open("assets/logo.png", "wb") as f:
                        f.write(uploaded_logo.getbuffer())
                    
                    st.success("Logo updated successfully! Refresh the page to see changes.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Error updating logo: {e}")

def initialize_session_state():
    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = []
    if 'knowledge_base_ready' not in st.session_state:
        st.session_state.knowledge_base_ready = check_knowledge_base()
    if 'query_history' not in st.session_state:
        st.session_state.query_history = []
    if 'selected_model' not in st.session_state:
        st.session_state.selected_model = "DeepSeek R1"
    if 'suggested_questions' not in st.session_state:
        st.session_state.suggested_questions = []
    if 'last_suggested_for' not in st.session_state:
        st.session_state.last_suggested_for = None
    if 'chat_history' not in st.session_state:
        # chat_history is a list of dicts: {"role": "user"|"assistant", "content": str}
        st.session_state.chat_history = []
    if 'authenticated' not in st.session_state:
        # Track whether the user is logged in
        st.session_state['authenticated'] = False
    if 'model_status' not in st.session_state:
        st.session_state.model_status = {}
    if 'last_model_check' not in st.session_state:
        st.session_state.last_model_check = 0
    if 'enhanced_prompt' not in st.session_state:
        st.session_state.enhanced_prompt = ""
    if 'enhance_prompt_error' not in st.session_state:
        st.session_state.enhance_prompt_error = ""

def check_knowledge_base():
    """Check if knowledge base files exist"""
    required_files = ["embeddings.npy", "faiss_index.index", "chunks.json"]
    return all(os.path.exists(f) for f in required_files)

def load_rag_components():
    """Load RAG system components"""
    try:
        embeddings = np.load("embeddings.npy")
        index = faiss.read_index("faiss_index.index")
        with open("chunks.json", 'r', encoding='utf-8') as f:
            chunks = json.load(f)
        # Use default embedding model from retreival.py
        embedding_model = EMBEDDING_MODELS["sentence-transformers"]
        return embeddings, index, chunks, embedding_model
    except Exception as e:
        st.error(f"Error loading RAG components: {e}")
        return None, None, None, None

def save_uploaded_files(uploaded_files):
    """Save uploaded files to data directory"""
    data_dir = "data"
    if not os.path.exists(data_dir):
        os.makedirs(data_dir)
    
    saved_files = []
    for uploaded_file in uploaded_files:
        try:
            file_path = os.path.join(data_dir, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            saved_files.append(uploaded_file.name)
        except Exception as e:
            st.error(f"Error saving {uploaded_file.name}: {e}")
    
    return saved_files

def get_file_stats():
    """Get statistics about processed files"""
    try:
        if os.path.exists("chunks.json"):
            with open("chunks.json", 'r', encoding='utf-8') as f:
                chunks = json.load(f)
            
            file_types = {}
            file_count = {}
            
            for chunk in chunks:
                filename = chunk.get('filename', 'Unknown')
                file_type = chunk.get('file_type', 'Unknown')
                
                if file_type not in file_types:
                    file_types[file_type] = 0
                file_types[file_type] += 1
                
                if filename not in file_count:
                    file_count[filename] = 0
                file_count[filename] += 1
            
            return len(chunks), len(file_count), file_types
        return 0, 0, {}
    except:
        return 0, 0, {}

def auto_initialize_knowledge_base():
    """Automatically initialize knowledge base if data directory has files but knowledge base doesn't exist"""
    data_dir = "data"
    required_files = ["embeddings.npy", "faiss_index.index", "chunks.json"]
    
    # Check if knowledge base files exist
    kb_exists = all(os.path.exists(f) for f in required_files)
    
    # Check if data directory has files
    data_has_files = False
    if os.path.exists(data_dir):
        supported_extensions = ['.pdf', '.xlsx', '.xls', '.csv', '.png', '.jpg', '.jpeg', '.tiff', '.bmp']
        files = [f for f in os.listdir(data_dir) if any(f.lower().endswith(ext) for ext in supported_extensions)]
        data_has_files = len(files) > 0
    
    # If data has files but knowledge base doesn't exist, create it
    if data_has_files and not kb_exists:
        st.info("Initializing knowledge base from existing files...")
        with st.spinner("Processing documents... This may take a few minutes."):
            try:
                from indexing import process_documents
                process_documents(data_dir, force_reprocess=True)
                st.success("Knowledge base initialized successfully!")
                st.session_state.knowledge_base_ready = True
                st.rerun()
            except Exception as e:
                st.error(f"Error initializing knowledge base: {e}")
                return False
    
    return kb_exists or data_has_files


def load_tabular_data(uploaded_file):
    """Load a CSV or Excel file into a pandas DataFrame."""
    if uploaded_file is None:
        return None

    name = uploaded_file.name.lower()
    try:
        if name.endswith(".csv"):
            try:
                df = pd.read_csv(uploaded_file)
            except UnicodeDecodeError:
                uploaded_file.seek(0)
                df = pd.read_csv(uploaded_file, encoding="latin1")
        elif name.endswith((".xlsx", ".xls")):
            df = pd.read_excel(uploaded_file)
        else:
            st.error("Unsupported file format. Please upload a CSV or Excel file.")
            return None
    except Exception as exc:
        st.error(f"Error reading file: {exc}")
        return None
    finally:
        try:
            uploaded_file.seek(0)
        except Exception:
            pass

    return df


def detect_datetime_and_numeric_columns(df: pd.DataFrame):
    """Detect datetime and numeric columns, coercing parseable date columns."""
    if df is None:
        return df, [], []

    converted_df = df.copy()
    datetime_cols = []

    for col in converted_df.columns:
        series = converted_df[col]
        if pd.api.types.is_datetime64_any_dtype(series):
            datetime_cols.append(col)
            continue

        if series.dtype == object:
            parsed = pd.to_datetime(series, errors="coerce")
            if parsed.notna().mean() >= 0.6:  # Treat as datetime when ≥60% parseable
                converted_df[col] = parsed
                datetime_cols.append(col)

    numeric_cols = converted_df.select_dtypes(include=[np.number]).columns.tolist()
    return converted_df, datetime_cols, numeric_cols


def infer_target_algorithms(df: pd.DataFrame, datetime_cols, numeric_cols):
    """Return compatibility table mapping numeric columns to available algorithms."""
    rows = []

    availability = {
        "SARIMAX": SARIMAX_AVAILABLE,
        "Prophet": PROPHET_AVAILABLE,
        "XGBoost": XGBOOST_AVAILABLE,
    }

    for target in numeric_cols:
        compat = []
        series = df[target].dropna()
        series_len = len(series)

        if not datetime_cols or series_len < 10:
            rows.append({
                "Target Column": target,
                "Compatible Models": "None (needs datetime column and ≥10 numeric values)",
            })
            continue

        if SARIMAX_AVAILABLE and series_len >= 24:
            compat.append("SARIMAX")
        if PROPHET_AVAILABLE and series_len >= 20:
            compat.append("Prophet")
        if XGBOOST_AVAILABLE and series_len >= 15:
            compat.append("XGBoost")

        if not compat:
            compat.append("None (not enough history)")

        rows.append({
            "Target Column": target,
            "Compatible Models": ", ".join(compat),
        })

    compatibility_df = pd.DataFrame(rows)

    unavailable = [model for model, available in availability.items() if not available]
    return compatibility_df, unavailable


def prepare_time_series(df: pd.DataFrame, date_col: str, target_col: str):
    """Prepare a time-series from dataframe columns."""
    ts_df = df[[date_col, target_col]].dropna()
    ts_df[date_col] = pd.to_datetime(ts_df[date_col], errors="coerce")
    ts_df = ts_df.dropna(subset=[date_col, target_col]).set_index(date_col).sort_index()
    ts_df.index = pd.DatetimeIndex(ts_df.index)

    if ts_df.empty:
        raise ValueError("No valid datetime/target pairs after cleaning.")

    freq = pd.infer_freq(ts_df.index)
    if not freq:
        # Fallback to daily frequency
        freq = "D"

    series = ts_df[target_col].astype(float)
    series = series.resample(freq).mean().ffill()
    return series, freq


def create_features(df: pd.DataFrame):
    """Create lag and calendar features for tree-based models."""
    feature_df = df.copy()
    feature_df.index = pd.DatetimeIndex(feature_df.index)
    feature_df["dayofweek"] = feature_df.index.dayofweek
    feature_df["quarter"] = feature_df.index.quarter
    feature_df["month"] = feature_df.index.month
    feature_df["year"] = feature_df.index.year
    feature_df["dayofyear"] = feature_df.index.dayofyear
    feature_df["lag1"] = feature_df.iloc[:, 0].shift(1)
    return feature_df.bfill()


def compute_confidence_series(forecast, lower_ci=None, upper_ci=None, resid_std=None):
    """Derive a pseudo-confidence score per forecasted point."""
    forecast_series = pd.Series(forecast)

    if lower_ci is not None and upper_ci is not None:
        lower_series = pd.Series(lower_ci)
        upper_series = pd.Series(upper_ci)
        interval = (upper_series - lower_series).abs()
        denom_series = forecast_series.abs() + interval + 1e-9

        interval_np = interval.to_numpy(dtype=float)
        denom_np = denom_series.to_numpy(dtype=float)
        ratio = np.divide(interval_np, denom_np, out=np.zeros_like(interval_np), where=denom_np != 0)
        confidence_vals = 1 - 0.5 * ratio
        confidence_series = pd.Series(confidence_vals, index=forecast_series.index)
        return confidence_series.clip(0.0, 0.99)

    if resid_std is not None:
        resid = float(abs(resid_std))
        reference = forecast_series.abs().mean()
        base = reference if reference else resid + 1e-6
        confidence = 1 - (resid / (base + resid + 1e-9))
        confidence = max(0.0, min(0.95, confidence))
        return pd.Series([confidence] * len(forecast_series), index=forecast_series.index)

    return pd.Series([np.nan] * len(forecast_series), index=forecast_series.index)


def get_sarimax_forecast(data: pd.Series, pred_length: int):
    if not SARIMAX_AVAILABLE or sm is None:
        raise ImportError("statsmodels is required for SARIMAX forecasting.")

    model = sm.tsa.SARIMAX(
        data,
        order=(1, 1, 1),
        seasonal_order=(1, 1, 1, 12),
        enforce_stationarity=False,
        enforce_invertibility=False,
    ).fit(disp=False)

    forecast_obj = model.get_forecast(steps=pred_length)
    forecast = forecast_obj.predicted_mean
    conf_int = forecast_obj.conf_int(alpha=0.05)
    resid_std = np.std(model.resid) if len(model.resid) else None

    return {
        "forecast": forecast,
        "lower_ci": conf_int.iloc[:, 0],
        "upper_ci": conf_int.iloc[:, 1],
        "confidence": compute_confidence_series(forecast, conf_int.iloc[:, 0], conf_int.iloc[:, 1], resid_std),
    }


def get_prophet_forecast(data: pd.Series, pred_length: int):
    if not PROPHET_AVAILABLE or Prophet is None:
        raise ImportError("prophet library is required for Prophet forecasting.")

    working = data.copy()
    working.index = pd.DatetimeIndex(working.index)

    prophet_df = working.reset_index().rename(columns={working.index.name or "ds": "ds", working.name or "value": "y"})
    prophet_df["ds"] = pd.to_datetime(prophet_df["ds"])

    model = Prophet(daily_seasonality=True)
    model.fit(prophet_df)

    date_index = pd.DatetimeIndex(working.index)
    freq = date_index.freq or pd.infer_freq(date_index)
    if not freq:
        freq = "D"
    future = model.make_future_dataframe(periods=pred_length, freq=freq)
    forecast = model.predict(future)
    future_forecast = forecast.iloc[-pred_length:][["ds", "yhat", "yhat_lower", "yhat_upper"]].set_index("ds")

    in_sample = model.predict(prophet_df)
    resid_std = np.std(prophet_df["y"] - in_sample["yhat"]) if len(prophet_df) else None

    return {
        "forecast": future_forecast["yhat"],
        "lower_ci": future_forecast["yhat_lower"],
        "upper_ci": future_forecast["yhat_upper"],
        "confidence": compute_confidence_series(
            future_forecast["yhat"],
            future_forecast["yhat_lower"],
            future_forecast["yhat_upper"],
            resid_std,
        ),
    }


def get_xgboost_forecast(data: pd.Series, pred_length: int):
    if not XGBOOST_AVAILABLE or xgb is None:
        raise ImportError("xgboost is required for XGBoost forecasting.")

    working = data.copy()
    working.index = pd.DatetimeIndex(working.index)

    df = pd.DataFrame(working)
    features = create_features(df.copy())
    target_name = working.name if working.name in features.columns else features.columns[0]
    target = features.pop(str(target_name))

    model = xgb.XGBRegressor(
        objective="reg:squarederror",
        n_estimators=800,
        learning_rate=0.03,
        max_depth=6,
        subsample=0.9,
        colsample_bytree=0.9,
        early_stopping_rounds=20,
        eval_metric="rmse",
    )

    model.fit(features, target, eval_set=[(features, target)], verbose=False)
    train_pred = model.predict(features)
    resid_std = np.std(target - train_pred) if len(train_pred) else None

    date_index = pd.DatetimeIndex(working.index)
    freq = date_index.freq or pd.infer_freq(date_index)
    if not freq:
        freq = "D"
    future_dates = pd.date_range(start=date_index.max(), periods=pred_length + 1, freq=freq)[1:]
    forecast_values = []
    current_lag = data.iloc[-1]

    for date in future_dates:
        temp_df = pd.DataFrame(index=[date])
        temp_df["dayofweek"] = date.dayofweek
        temp_df["quarter"] = date.quarter
        temp_df["month"] = date.month
        temp_df["year"] = date.year
        temp_df["dayofyear"] = date.dayofyear
        temp_df["lag1"] = current_lag

        prediction = model.predict(temp_df)[0]
        forecast_values.append(prediction)
        current_lag = prediction

    forecast_series = pd.Series(forecast_values, index=future_dates)
    return {
        "forecast": forecast_series,
        "lower_ci": None,
        "upper_ci": None,
        "confidence": compute_confidence_series(forecast_series, resid_std=resid_std),
    }


def backtest_model(data: pd.Series, model_name: str, test_size: int):
    train = data.iloc[:-test_size]
    test = data.iloc[-test_size:]

    try:
        if model_name == "SARIMAX":
            preds = get_sarimax_forecast(train, len(test))
        elif model_name == "Prophet":
            preds = get_prophet_forecast(train, len(test))
        elif model_name == "XGBoost":
            preds = get_xgboost_forecast(train, len(test))
        else:
            return None

        mape = mean_absolute_percentage_error(test, preds["forecast"]) * 100
        return f"{mape:.2f}% Error"
    except Exception as exc:
        return f"Failed: {exc}"


def stream_llm_response(prompt: str, temperature: float = 0.5, max_tokens: int = 700):
    """Call the LLM with streaming response and aggregate output."""
    if not OPENAI_TOKEN:
        return "LLM token not configured."

    try:
        client = create_llm_client()
        response = client.chat.completions.create(
            model=OPENAI_DEPLOYMENT,
            messages=[{"role": "user", "content": prompt}],
            temperature=temperature,
            max_tokens=max_tokens,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=1,
            stream=True,
        )

        collected = ""
        for chunk in response:
            if hasattr(chunk, "choices") and chunk.choices:
                delta = chunk.choices[0].delta
                if delta and getattr(delta, "content", None):
                    collected += delta.content
        return collected.strip()
    except Exception as exc:
        return f"Error calling LLM: {exc}"


def enhance_query_prompt(prompt: str, selected_model: Optional[str]) -> Tuple[bool, str]:
    """Generate a refined version of the user's question using the selected LLM."""
    if not prompt.strip():
        return False, "Enter a question before enhancing the prompt."

    if not selected_model:
        return False, "Select an AI model before enhancing the prompt."

    model_config = MODELS.get(selected_model)
    if not model_config:
        return False, f"Model '{selected_model}' is not configured."

    try:
        client = create_client_for_model(selected_model)
    except Exception as exc:  # pragma: no cover - surface upstream errors to the user
        return False, f"Could not connect to {selected_model}: {exc}"

    enhancement_instructions = (
        "You are a careful editor. Polish the question with minimal grammatical improvements or gentle reordering only. "
        "Do not change the meaning, do not add or remove context, and do not introduce new facts or instructions. "
        "Return a single revised prompt that preserves the user's intent exactly."
    )

    payload = [
        {"role": "system", "content": enhancement_instructions},
        {"role": "user", "content": prompt.strip()},
    ]

    try:
        completion = client.chat.completions.create(
            model=model_config["deployment"],
            messages=payload,  # type: ignore[arg-type]
            max_tokens=120,
            temperature=0.2,
            top_p=0.6,
        )
        choice = completion.choices[0]
        content = getattr(choice.message, "content", "") or ""
        enhanced_text = content.strip()
        if not enhanced_text:
            return False, "The model returned an empty prompt. Please try again."
        return True, enhanced_text
    except Exception as exc:
        return False, f"Error enhancing prompt: {exc}"


def generate_dataset_summary(df: pd.DataFrame, file_name: str):
    head_df = df.head(5)
    try:
        preview = head_df.to_markdown(index=False)
    except Exception:
        preview = head_df.to_string(index=False)

    try:
        stats_df = df.describe(include="all").transpose().head(10)
    except Exception:
        stats_df = None

    if stats_df is not None:
        try:
            description = stats_df.to_markdown()
        except Exception:
            description = stats_df.to_string()
    else:
        description = "Summary unavailable."
    prompt = f"""
You are a senior data analyst. Provide a concise bullet summary of the uploaded dataset.

File name: {file_name}
Shape: {df.shape[0]} rows × {df.shape[1]} columns

Preview (first 5 rows):
{preview}

Key statistics:
{description}

Highlight:
1. Notable distributions or anomalies.
2. Columns that look like potential time series targets.
3. Data quality concerns to fix before forecasting.

Keep the answer short, use bullet points, and address non-technical users.
"""
    return stream_llm_response(prompt, temperature=0.4, max_tokens=800)


def generate_forecast_insight(date_col, target_col, model_errors, forecast_table):
    forecast_markdown = forecast_table.reset_index().head(12).to_markdown(index=False)
    error_lines = "\n".join([f"- {model}: {error}" for model, error in model_errors.items()]) or "- No backtest available"
    prompt = f"""
You are an analytics co-pilot. Summarize the forecasting results for stakeholders.

Date column: {date_col}
Target column: {target_col}

Model accuracy snapshot (MAPE):
{error_lines}

Forecast preview:
{forecast_markdown}

Explain:
1. Overall trend in the forecast horizon.
2. Which model performed best and why.
3. Recommended next steps.

Keep the tone friendly and concise. Use short paragraphs or bullet points.
"""
    return stream_llm_response(prompt, temperature=0.45, max_tokens=700)


def add_table_slide(presentation: Any, title: str, dataframe: pd.DataFrame):
    if not PPTX_AVAILABLE or Presentation is None or Inches is None or Pt is None:
        return

    slide_layout = presentation.slide_layouts[5]  # Title and content
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    rows, cols = dataframe.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    table = table_shape.table

    for idx, column in enumerate(dataframe.columns):
        header_cell = table.cell(0, idx)
        header_cell.text = str(column)
        header_cell.text_frame.paragraphs[0].font.bold = True
        header_cell.text_frame.paragraphs[0].font.size = Pt(12)

    for row_idx in range(rows):
        for col_idx in range(cols):
            value = dataframe.iat[row_idx, col_idx]
            data_cell = table.cell(row_idx + 1, col_idx)
            data_cell.text = "" if pd.isna(value) else str(value)
            data_cell.text_frame.paragraphs[0].font.size = Pt(11)


def add_line_chart_slide(presentation: Any, title: str, chart_df: pd.DataFrame):
    if (
        not PPTX_AVAILABLE
        or Presentation is None
        or CategoryChartData is None
        or XL_CHART_TYPE is None
        or Inches is None
    ):
        return

    if chart_df is None or chart_df.empty or "Date" not in chart_df.columns:
        return

    value_columns = [col for col in ["Actual", "Forecast"] if col in chart_df.columns]
    if not value_columns:
        return

    plot_df = chart_df.copy().sort_values("Date").tail(60)

    categories = []
    for date_val in plot_df["Date"]:
        if isinstance(date_val, pd.Timestamp):
            categories.append(date_val.strftime("%Y-%m-%d"))
        elif isinstance(date_val, datetime.date):
            categories.append(date_val.strftime("%Y-%m-%d"))
        else:
            categories.append(str(date_val))

    chart_data = CategoryChartData()
    chart_data.categories = categories

    for column in value_columns:
        numeric_series = pd.to_numeric(plot_df[column], errors="coerce")
        values = [None if pd.isna(val) else float(val) for val in numeric_series]
        chart_data.add_series(column, values)

    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(4.5),
        chart_data,
    )

    chart = chart_shape.chart
    chart.has_legend = True


def build_ppt_report(
    title: str,
    dataset_summary: str,
    model_errors: dict,
    forecast_table: pd.DataFrame,
    insights: str,
    forecast_chart_df: Optional[pd.DataFrame] = None,
):
    if not PPTX_AVAILABLE or Presentation is None:
        return None

    presentation = Presentation()

    # Title slide
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title  # type: ignore[attr-defined]
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = (dataset_summary or "Forecast overview generated by AI.")[:500]  # type: ignore[attr-defined]

    # Model accuracy slide
    if model_errors:
        errors_df = pd.DataFrame([
            {"Model": model, "MAPE": error} for model, error in model_errors.items()
        ])
        add_table_slide(presentation, "Model Accuracy", errors_df)

    # Forecast vs history chart
    if forecast_chart_df is not None and not forecast_chart_df.empty:
        add_line_chart_slide(presentation, "Forecast vs History", forecast_chart_df)

    # Forecast table slides (paged for readability)
    if forecast_table is not None and not forecast_table.empty:
        forecast_for_ppt = forecast_table.reset_index().copy()
        forecast_for_ppt.columns = ["Date"] + list(forecast_for_ppt.columns[1:])

        chunk_size = 15
        total_rows = len(forecast_for_ppt)

        for chunk_idx, start in enumerate(range(0, total_rows, chunk_size)):
            chunk = forecast_for_ppt.iloc[start:start + chunk_size]
            if chunk.empty:
                continue

            end = start + len(chunk)
            if chunk_idx == 0:
                slide_title = "Forecast Outlook"
            else:
                slide_title = f"Forecast Table ({start + 1}-{end})"

            add_table_slide(presentation, slide_title, chunk)

    # Insight slide
    insight_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    insight_slide.shapes.title.text = "Key Insights"  # type: ignore[attr-defined]
    textbox = insight_slide.shapes.placeholders[1]
    textbox.text = insights or "Insights will appear here once generated."  # type: ignore[attr-defined]

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def render_forecasting_workspace():
    st.title("📈 Forecasting Workspace")
    st.caption("Upload tabular data to explore trends, run forecasts, and export insights.")

    uploaded = st.file_uploader(
        "Upload CSV or Excel",
        type=["csv", "xlsx", "xls"],
        key="forecast_main_uploader",
    )

    def _reset_forecast_state():
        for key in [
            "forecast_original_df",
            "forecast_df",
            "forecast_datetime_cols",
            "forecast_numeric_cols",
            "forecast_file_name",
            "forecast_model_errors",
            "forecast_result_table",
            "forecast_selected_model",
            "forecast_chart_df",
            "forecast_llm_analysis",
            "forecast_llm_insight",
            "forecast_explain_llm",
            "forecast_ppt_bytes",
            "forecast_series_freq",
            "forecast_results",
            "forecast_errors",
            "forecast_charts",
            "forecast_selected_targets",
        ]:
            st.session_state.pop(key, None)

    if uploaded is not None:
        df = load_tabular_data(uploaded)
        if df is not None and not df.empty:
            _reset_forecast_state()
            converted_df, datetime_cols, numeric_cols = detect_datetime_and_numeric_columns(df)
            st.session_state["forecast_original_df"] = df
            st.session_state["forecast_df"] = converted_df
            st.session_state["forecast_datetime_cols"] = datetime_cols
            st.session_state["forecast_numeric_cols"] = numeric_cols
            st.session_state["forecast_file_name"] = uploaded.name
            st.success(f"Loaded {uploaded.name}")
        else:
            st.error("Uploaded file appears to be empty.")

    df = st.session_state.get("forecast_df")
    datetime_cols = st.session_state.get("forecast_datetime_cols", [])
    numeric_cols = st.session_state.get("forecast_numeric_cols", [])
    file_name = st.session_state.get("forecast_file_name", "uploaded file")

    if df is None or df.empty:
        st.info("Upload a CSV or Excel file to begin forecasting.")
        return

    st.markdown("**Data Preview**")
    st.dataframe(df.head(), use_container_width=True)

    compatibility_df, unavailable_models = infer_target_algorithms(df, datetime_cols, numeric_cols)
    if unavailable_models:
        st.warning(
            "Unavailable model libraries: " + ", ".join(unavailable_models) + ". Install them to unlock more forecasts."
        )

    if not compatibility_df.empty:
        st.markdown("**Model Compatibility**")
        st.dataframe(compatibility_df, use_container_width=True)

    st.markdown("**LLM Data Scan**")
    if st.button("Analyze Dataset with LLM", key="forecast_llm_button"):
        with st.spinner("Summarizing dataset with LLM..."):
            summary = generate_dataset_summary(df, file_name)
            st.session_state["forecast_llm_analysis"] = summary

    if st.session_state.get("forecast_llm_analysis"):
        st.markdown(st.session_state["forecast_llm_analysis"])

    if not datetime_cols:
        st.warning("No datetime columns detected. Convert a column to datetime format to enable forecasting.")
        return

    forecast_map = {}
    for _, row in compatibility_df.iterrows():
        target = row["Target Column"]
        models = [m.strip() for m in str(row["Compatible Models"]).split(",") if "None" not in m]
        forecast_map[target] = [m for m in models if m in {"SARIMAX", "Prophet", "XGBoost"}]

    date_col = st.selectbox(
        "Select date column",
        options=datetime_cols,
        index=0,
        key="forecast_date_col",
    )

    target_options = [col for col, models in forecast_map.items() if models]
    if not target_options:
        st.warning("No numeric columns qualify for forecasting. Ensure the dataset has at least 15 numeric entries.")
        return

    default_targets = target_options[:1]
    selected_targets = st.multiselect(
        "Select target columns",
        options=target_options,
        default=default_targets,
        key="forecast_target_cols",
        help="Choose one or more numeric columns to forecast in a single run.",
    )

    if not selected_targets:
        st.warning("Select at least one target column to continue.")
        return

    st.session_state["forecast_selected_targets"] = selected_targets

    preview_target = st.selectbox(
        "Preview target for visuals",
        options=selected_targets,
        index=0,
        key="forecast_preview_target",
        help="Used for the quick chart below; all selected targets are forecasted.",
    )

    # Determine models that are supported for every selected target
    common_models = set(forecast_map[selected_targets[0]])
    for tgt in selected_targets[1:]:
        common_models &= set(forecast_map.get(tgt, []))

    if not common_models:
        st.error("No common forecasting models are available for the selected columns. Try removing a column or installing the required libraries.")
        return

    available_models = sorted(common_models)

    selected_model = st.selectbox(
        "Choose forecasting model",
        options=available_models,
        key="forecast_model_choice",
    )

    horizon = st.number_input(
        "Forecast horizon (future steps)",
        min_value=1,
        max_value=365,
        value=30,
        key="forecast_horizon",
    )

    st.markdown("**Data Visualizations**")
    chart_type = st.selectbox(
        "Choose visualization",
        options=["Line Trend", "Distribution"],
        key="forecast_chart_type",
    )
    try:
        trend_df = df[[date_col, preview_target]].dropna().sort_values(date_col)
        trend_df = trend_df.tail(min(len(trend_df), 200))
        trend_df = trend_df.rename(columns={date_col: "Date", preview_target: "Value"})

        if chart_type == "Line Trend":
            fig = px.line(trend_df, x="Date", y="Value", title=f"{preview_target} over time")
            fig.update_layout(height=320)
            st.plotly_chart(fig, use_container_width=True)
        else:
            dist_df = pd.DataFrame({preview_target: df[preview_target].dropna()})
            fig = px.histogram(dist_df, x=preview_target, nbins=30, title=f"Distribution of {preview_target}")
            fig.update_layout(height=320)
            st.plotly_chart(fig, use_container_width=True)
    except Exception as exc:
        st.warning(f"Unable to plot data visualization: {exc}")

    run_forecast = st.button("Run Forecast", type="primary", key="run_forecast_button")

    if run_forecast:
        if not date_col or not selected_targets:
            st.warning("Please select both date and target columns before forecasting.")
            return

        # Clean up legacy single-target state
        for legacy_key in ["forecast_result_table", "forecast_chart_df", "forecast_model_errors"]:
            st.session_state.pop(legacy_key, None)

        results: Dict[str, pd.DataFrame] = {}
        errors: Dict[str, Dict[str, str]] = {}
        charts: Dict[str, pd.DataFrame] = {}
        insights: Dict[str, str] = {}
        narratives: Dict[str, str] = {}
        ppt_buffers: Dict[str, bytes] = {}
        freq_map: Dict[str, str] = {}
        failed_targets: List[str] = []

        total_targets = len(selected_targets)

        for idx, target_col in enumerate(selected_targets, start=1):
            try:
                with st.spinner(f"Running forecast for {target_col} ({idx}/{total_targets})..."):
                    series, freq = prepare_time_series(df, date_col, target_col)
                    freq_map[target_col] = freq

                    model_errors: Dict[str, str] = {}
                    test_size = max(1, min(len(series) // 5, 30))
                    if len(series) > test_size * 2:
                        for model_name in available_models:
                            error = backtest_model(series, model_name, test_size)
                            if error:
                                model_errors[model_name] = error
                    else:
                        st.info(f"{target_col}: Not enough history for backtesting. Skipping MAPE evaluation.")

                    forecast_outputs = {}
                    for model_name in available_models:
                        try:
                            if model_name == "SARIMAX":
                                forecast_outputs[model_name] = get_sarimax_forecast(series, horizon)
                            elif model_name == "Prophet":
                                forecast_outputs[model_name] = get_prophet_forecast(series, horizon)
                            elif model_name == "XGBoost":
                                forecast_outputs[model_name] = get_xgboost_forecast(series, horizon)
                        except Exception as exc:
                            st.warning(f"{target_col} - {model_name} failed: {exc}")

                    if selected_model not in forecast_outputs:
                        st.error(
                            f"{target_col}: Selected model did not return results. Choose a different model or adjust data."
                        )
                        failed_targets.append(target_col)
                        continue

                    chosen = forecast_outputs[selected_model]
                    forecast_series = chosen["forecast"]
                    lower_ci = chosen.get("lower_ci")
                    upper_ci = chosen.get("upper_ci")
                    confidence_series = chosen.get("confidence")

                    if lower_ci is None:
                        lower_ci = pd.Series([np.nan] * len(forecast_series), index=forecast_series.index)
                    if upper_ci is None:
                        upper_ci = pd.Series([np.nan] * len(forecast_series), index=forecast_series.index)
                    if confidence_series is None:
                        confidence_series = pd.Series([np.nan] * len(forecast_series), index=forecast_series.index)

                    result_table = pd.DataFrame({
                        "Forecast": forecast_series,
                        "Lower 95%": lower_ci,
                        "Upper 95%": upper_ci,
                        "Confidence Score": confidence_series,
                    })
                    result_table.index.name = "Date"
                    result_table["Forecast"] = pd.to_numeric(result_table["Forecast"], errors="coerce")
                    result_table["Lower 95%"] = pd.to_numeric(result_table["Lower 95%"], errors="coerce")
                    result_table["Upper 95%"] = pd.to_numeric(result_table["Upper 95%"], errors="coerce")
                    result_table["Confidence Score"] = pd.to_numeric(result_table["Confidence Score"], errors="coerce").clip(0.0, 0.99)

                    chart_history = series.tail(min(len(series), 200)).to_frame("Actual")
                    chart_forecast = chosen["forecast"].to_frame("Forecast")
                    chart_combined = pd.concat([chart_history, chart_forecast])
                    chart_combined = chart_combined.reset_index().rename(columns={"index": "Date"})

                    insights[target_col] = generate_forecast_insight(date_col, target_col, model_errors, result_table)
                    narratives[target_col] = get_forecasting_explanation(
                        all_columns=list(df.columns),
                        date_col=date_col,
                        target_col=target_col,
                        model_errors=model_errors or {selected_model: "Not evaluated"},
                    )

                    ppt_buffer = build_ppt_report(
                        title=f"Forecast Report – {target_col}",
                        dataset_summary=st.session_state.get("forecast_llm_analysis", ""),
                        model_errors=model_errors,
                        forecast_table=result_table,
                        insights=insights[target_col],
                        forecast_chart_df=chart_combined,
                    )
                    if ppt_buffer is not None:
                        ppt_buffers[target_col] = ppt_buffer

                    results[target_col] = result_table
                    errors[target_col] = model_errors
                    charts[target_col] = chart_combined
            except Exception as exc:
                st.error(f"{target_col}: Forecasting error: {exc}")
                failed_targets.append(target_col)

        if not results:
            st.error("Forecasting failed for all selected targets. Review errors above and try again.")
            return

        st.session_state["forecast_results"] = results
        st.session_state["forecast_errors"] = errors
        st.session_state["forecast_charts"] = charts
        st.session_state["forecast_llm_insight"] = insights
        st.session_state["forecast_explain_llm"] = narratives
        st.session_state["forecast_ppt_bytes"] = ppt_buffers
        st.session_state["forecast_selected_model"] = selected_model
        st.session_state["forecast_series_freq"] = freq_map

        success_msg = f"Forecast complete for {len(results)} target(s)."
        if failed_targets:
            success_msg += " Skipped: " + ", ".join(failed_targets)
        st.success(success_msg)

    forecast_results = st.session_state.get("forecast_results", {})
    if forecast_results:
        st.markdown("**Forecast Output**")
        target_tabs = st.tabs(list(forecast_results.keys()))

        charts_by_target = st.session_state.get("forecast_charts", {})
        errors_by_target = st.session_state.get("forecast_errors", {})
        insights_by_target = st.session_state.get("forecast_llm_insight", {})
        narratives_by_target = st.session_state.get("forecast_explain_llm", {})
        ppt_by_target = st.session_state.get("forecast_ppt_bytes", {})

        for tab, target_col in zip(target_tabs, forecast_results.keys()):
            with tab:
                display_df = forecast_results[target_col].copy()
                st.dataframe(
                    display_df.style.format({
                        "Forecast": "{:,.2f}",
                        "Lower 95%": "{:,.2f}",
                        "Upper 95%": "{:,.2f}",
                        "Confidence Score": "{:.2%}",
                    }),
                    use_container_width=True,
                )

                chart_df = charts_by_target.get(target_col)
                if chart_df is not None and not chart_df.empty:
                    st.markdown("**Forecast vs History**")
                    fig = px.line(
                        chart_df,
                        x="Date",
                        y=[col for col in ["Actual", "Forecast"] if col in chart_df.columns],
                        title=f"Forecast vs History – {target_col}",
                    )
                    fig.update_layout(height=320)
                    st.plotly_chart(fig, use_container_width=True)

                model_errors = errors_by_target.get(target_col, {})
                if model_errors:
                    st.markdown("**Model Accuracy (MAPE)**")
                    for model, error in model_errors.items():
                        st.write(f"- {model}: {error}")

                insight_text = insights_by_target.get(target_col)
                if insight_text:
                    st.markdown("**AI Insight**")
                    st.markdown(insight_text)

                narrative_text = narratives_by_target.get(target_col)
                if narrative_text:
                    st.markdown("**Narrative Recommendation**")
                    st.markdown(narrative_text)

                ppt_bytes = ppt_by_target.get(target_col)
                if ppt_bytes is not None:
                    st.download_button(
                        f"Download PPT Report ({target_col})",
                        data=ppt_bytes,
                        file_name=f"forecast_report_{target_col}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"forecast_ppt_download_{target_col}",
                    )

        if not ppt_by_target:
            st.info("Install python-pptx to enable PPT downloads.")


def render_sidebar_kb_controls():
    st.subheader("Knowledge Base Status")
    if st.session_state.knowledge_base_ready:
        chunks_count, files_count, file_types = get_file_stats()
        st.success(f"Ready - {files_count} files, {chunks_count} chunks")

        if file_types:
            st.write("File Types:")
            for file_type, count in file_types.items():
                st.write(f"  • {file_type}: {count} chunks")
    else:
        st.warning("No knowledge base found. Upload and process files first.")

    st.divider()
    st.subheader("Actions")

    if st.button("Refresh Knowledge Base", help="Check for new files and update"):
        with st.spinner("Refreshing..."):
            try:
                process_documents()
                st.session_state.knowledge_base_ready = check_knowledge_base()
                st.success("Knowledge base refreshed!")
                st.rerun()
            except Exception as e:
                st.error(f"Error refreshing: {e}")

    if st.button("Force Rebuild", help="Rebuild entire knowledge base from scratch"):
        with st.spinner("Clearing existing data and rebuilding..."):
            try:
                files_to_clear = ["embeddings.npy", "faiss_index.index", "chunks.json"]
                for file_path in files_to_clear:
                    if os.path.exists(file_path):
                        os.remove(file_path)

                process_documents(force_reprocess=True)
                st.session_state.knowledge_base_ready = check_knowledge_base()

                if st.session_state.knowledge_base_ready:
                    st.success("Knowledge base completely rebuilt from scratch!")
                else:
                    st.warning("Rebuild completed but knowledge base files not found. Please ensure you have documents in the data directory.")

                st.rerun()
            except Exception as e:
                st.error(f"Error rebuilding: {e}")
                st.session_state.knowledge_base_ready = check_knowledge_base()

    st.divider()
    st.subheader("📁 File Management")
    processed_files = get_processed_files()

    if processed_files:
        st.write(f"**Processed Files ({len(processed_files)}):**")

        for file_info in processed_files:
            filename = file_info["filename"]
            chunks = file_info["chunks"]

            col_file, col_delete = st.columns([3, 1])

            with col_file:
                st.write(f"📄 {filename}")
                st.caption(f"{chunks} chunks")

            with col_delete:
                if st.button("🗑️", key=f"delete_{filename}", help=f"Delete {filename}"):
                    if st.session_state.get(f'confirm_delete_{filename}', False):
                        with st.spinner(f"Deleting {filename}..."):
                            try:
                                success = delete_file_from_knowledge_base(filename)
                                if success:
                                    st.success(f"Deleted {filename}")
                                    st.session_state.knowledge_base_ready = check_knowledge_base()
                                    st.session_state[f'confirm_delete_{filename}'] = False
                                    st.rerun()
                                else:
                                    st.error(f"Failed to delete {filename}")
                            except Exception as e:
                                st.error(f"Error deleting {filename}: {e}")
                                st.session_state[f'confirm_delete_{filename}'] = False
                    else:
                        st.session_state[f'confirm_delete_{filename}'] = True
                        st.warning(f"⚠️ Click again to confirm deletion of {filename}")
                        st.rerun()
    else:
        st.info("No processed files found.")

    st.divider()

    if st.button("🗑️ Clear All Data", help="Remove all processed data and uploaded files", type="secondary"):
        if st.session_state.get('confirm_clear', False):
            with st.spinner("Clearing all data..."):
                try:
                    files_to_clear = ["embeddings.npy", "faiss_index.index", "chunks.json"]
                    for file_path in files_to_clear:
                        if os.path.exists(file_path):
                            os.remove(file_path)

                    if os.path.exists("data"):
                        shutil.rmtree("data")
                        os.makedirs("data")

                    st.session_state.knowledge_base_ready = False
                    st.session_state.processed_files = []
                    st.session_state.query_history = []
                    st.session_state.confirm_clear = False

                    st.success("All data cleared successfully!")
                    st.rerun()
                except Exception as e:
                    st.error(f"Error clearing data: {e}")
                    st.session_state.confirm_clear = False
        else:
            st.session_state.confirm_clear = True
            st.warning("⚠️ Click again to confirm - this will delete ALL data!")
            st.rerun()

def main():
    load_css()
    initialize_session_state()
    
    # If not authenticated, show login form and return early
    if not st.session_state.get('authenticated', False):
        st.title("Please log in")
        with st.form("login_form"):
            uname = st.text_input("Username")
            pwd = st.text_input("Password", type="password")
            submit = st.form_submit_button("Login")
            if submit:
                if CREDENTIALS.get(uname) == pwd:
                    st.session_state['authenticated'] = True
                    st.success(f"Logged in as {uname}")
                    st.rerun()
                else:
                    st.error("Invalid username or password")
        return

    # Auto-initialize knowledge base if needed
    auto_initialize_knowledge_base()
    
    # Header
    logo_data = load_logo()
    
    if logo_data:
        # Header with logo
        st.markdown(f"""
        <div class="header-container">
            <div class="logo-container">
                <img src="{logo_data}" class="company-logo" alt="Company Logo">
                <h1 class="title">Multimodal GenAI Chatbot System</h1>
            </div>
            <div class="features-container">
                <div class="feature-item"> Intelligent Document Processing</div>
                <div class="feature-item"> Question Answering using RAG (Retrieval Augmented Generation)</div>
                <div class="feature-item"> Seamless Integration with Website</div>
                <div class="feature-item"> Secure Access Control for Department Admins</div>
            </div>
        </div>
        """, unsafe_allow_html=True)
    else:
        # Header without logo (fallback)
        st.markdown("""
        <div class="header-container">
            <div class="logo-container">
                <h1 class="title">Multimodal GenAI Chatbot System</h1>
            </div>
            <div class="features-container">
                <div class="feature-item"> Intelligent Document Processing</div>
                <div class="feature-item"> Question Answering using RAG (Retrieval Augmented Generation)</div>
                <div class="feature-item"> Seamless Integration with Website</div>
                <div class="feature-item"> Secure Access Control for Department Admins</div>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    workspace_choice_default = st.session_state.get("last_workspace", "Knowledge Base")
    workspace_choice = workspace_choice_default

    # Sidebar
    with st.sidebar:
        st.header("Navigation")
        workspace_choice = st.radio(
            "Select workspace",
            options=["Knowledge Base", "Forecasting"],
            index=0 if workspace_choice_default == "Knowledge Base" else 1,
        )
        st.session_state["last_workspace"] = workspace_choice

        st.divider()
        if workspace_choice == "Knowledge Base":
            render_sidebar_kb_controls()
        else:
            st.info("Use the Forecasting workspace to upload data and run models from the main panel.")

        st.divider()
        upload_logo_interface()

        st.divider()
        if st.button("Logout", key="sidebar_logout"):
            st.session_state['authenticated'] = False
            st.success("Logged out")
            st.rerun()

    if workspace_choice == "Forecasting":
        render_forecasting_workspace()
        return

    # Main content area with columns
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.markdown('<div class="upload-section">', unsafe_allow_html=True)
        st.subheader("📁 Upload Documents")
        
        uploaded_files = st.file_uploader(
            "Choose files",
            type=['pdf', 'xlsx', 'xls', 'csv', 'png', 'jpg', 'jpeg', 'tiff', 'bmp'],
            accept_multiple_files=True,
            help="Supported formats: PDF, Excel, CSV, Images"
        )
        
        if uploaded_files:
            st.write(f"Selected {len(uploaded_files)} file(s):")
            for file in uploaded_files:
                st.write(f"  • {file.name} ({file.size} bytes)")
            
            if st.button("Process Files", type="primary"):
                with st.spinner("Processing files..."):
                    try:
                        # Save uploaded files
                        saved_files = save_uploaded_files(uploaded_files)
                        
                        if saved_files:
                            st.success(f"Saved {len(saved_files)} files")
                            
                            # Process documents
                            process_documents()
                            st.session_state.knowledge_base_ready = check_knowledge_base()
                            st.success("Files processed successfully!")
                            st.session_state.processed_files.extend(saved_files)

                            # Load newly processed chunks and extract suggested questions for UI
                            try:
                                with open('chunks.json', 'r', encoding='utf-8') as f:
                                    chunks = json.load(f)

                                # Find suggestions attached to documents we just saved
                                new_suggestions = []
                                for ch in chunks:
                                    fn = ch.get('filename')
                                    if fn in saved_files and ch.get('suggested_questions'):
                                        for q in ch.get('suggested_questions'):
                                            if q not in new_suggestions:
                                                new_suggestions.append(q)

                                if new_suggestions:
                                    st.session_state.suggested_questions = new_suggestions
                                    st.session_state.last_suggested_for = saved_files
                            except Exception as e:
                                st.warning(f"Could not load suggested questions: {e}")

                            st.rerun()
                        else:
                            st.error("No files were saved")
                            
                    except Exception as e:
                        st.error(f"Error processing files: {e}")
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col2:
        st.markdown('<div class="query-section">', unsafe_allow_html=True)
        st.subheader("❓ Ask Questions")
        
        if st.session_state.knowledge_base_ready:
            # Check model status periodically
            import time
            current_time = time.time()
            if current_time - st.session_state.last_model_check > 300:  # Check every 5 minutes
                with st.spinner("Checking model availability..."):
                    st.session_state.model_status = get_available_models()
                    st.session_state.last_model_check = current_time
            
            # Model selection
            st.subheader("🤖 Select AI Model")
            
            available_models = []
            offline_models = []
            
            for model, status in st.session_state.model_status.items():
                if status:
                    available_models.append(f"{model} 🟢")
                else:
                    offline_models.append(f"{model} 🔴")
            
            all_model_options = available_models + offline_models
            
            if available_models:
                # Default selection
                if st.session_state.selected_model not in [m.replace(" 🟢", "").replace(" 🔴", "") for m in all_model_options]:
                    st.session_state.selected_model = available_models[0].replace(" 🟢", "")
                
                # Find current selection index
                current_model_display = f"{st.session_state.selected_model} 🟢" if st.session_state.selected_model in [m.replace(" 🟢", "") for m in available_models] else f"{st.session_state.selected_model} 🔴"
                try:
                    default_index = all_model_options.index(current_model_display)
                except ValueError:
                    default_index = 0
                
                selected_model_display = st.selectbox(
                    "Choose AI model:",
                    options=all_model_options,
                    index=default_index,
                    help="🟢 = Online, 🔴 = Offline"
                )
                
                if selected_model_display:
                    st.session_state.selected_model = selected_model_display.replace(" 🟢", "").replace(" 🔴", "")
                
                # Show model status
                if selected_model_display and "🔴" in selected_model_display:
                    st.warning(f"⚠️ {st.session_state.selected_model} is currently offline. Please select an online model.")
                else:
                    st.success(f"✅ {st.session_state.selected_model} is ready")
            else:
                st.error("❌ No models are currently online. Please try again later.")
                st.session_state.selected_model = None
            
            # Manual refresh button
            col_refresh, col_space = st.columns([1, 3])
            with col_refresh:
                if st.button("🔄 Refresh Models"):
                    with st.spinner("Checking models..."):
                        st.session_state.model_status = get_available_models()
                        st.session_state.last_model_check = current_time
                        st.rerun()
            
            # Embedding Model Selection
            st.subheader("🏛️ Embedding Model")
            
            # Initialize embedding model selection
            if 'embedding_model_choice' not in st.session_state:
                st.session_state.embedding_model_choice = "Auto-detect"
            
            embedding_options = [
                "Auto-detect (Smart legal/general selection)",
                "Sentence Transformers (General purpose)",
                "InLegalBERT (Legal documents)"
            ]
            
            selected_embedding = st.selectbox(
                "Choose embedding model:",
                options=embedding_options,
                index=embedding_options.index(st.session_state.embedding_model_choice) if st.session_state.embedding_model_choice in embedding_options else 0,
                help="Auto-detect will choose InLegalBERT for legal content and Sentence Transformers for general content"
            )
            
            st.session_state.embedding_model_choice = selected_embedding
            
            if "InLegalBERT" in selected_embedding:
                st.info("🏛️ InLegalBERT provides specialized embeddings for legal documents, contracts, and case law analysis.")
            elif "Auto-detect" in selected_embedding:
                st.info("🔍 Smart detection will automatically choose the best embedding model based on your query and document content.")
            else:
                st.info("📚 Sentence Transformers provides high-quality general-purpose embeddings for all types of documents.")
            
            st.divider()
            
            # Load RAG components
            embeddings, index, chunks, embedding_model = load_rag_components()
            
            if all(component is not None for component in [embeddings, index, chunks, embedding_model]):
                # If an autofill query exists (from clicking a suggestion), use it
                prefill = st.session_state.pop('_autofill_query', '') if '_autofill_query' in st.session_state else ''
                query = st.text_area(
                    "Enter your question:",
                    value=prefill,
                    height=120,
                    placeholder="Ask anything about your uploaded documents...",
                )

                # If there are suggested questions from the latest upload, add them as an assistant bubble
                # so they appear in the chat history (once per upload batch).
                if st.session_state.suggested_questions:
                    # Use last_suggested_for to avoid adding the same assistant bubble repeatedly
                    flag_key = '_suggested_shown_for'
                    last_for = st.session_state.get('last_suggested_for')
                    already_shown_for = st.session_state.get(flag_key)

                    if last_for and already_shown_for != last_for:
                        # Create an assistant message with followups equal to suggested questions
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": "Suggested questions from recently uploaded documents:",
                            "followups": st.session_state.suggested_questions
                        })
                        st.session_state[flag_key] = last_for
                
                # Only allow questions if model is online
                model_online = bool(
                    st.session_state.selected_model
                    and st.session_state.model_status.get(st.session_state.selected_model, False)
                )

                action_cols = st.columns(2)
                with action_cols[0]:
                    get_answer_clicked = st.button(
                        "Get Answer",
                        type="primary",
                        key="btn_get_answer",
                        disabled=not query.strip() or not model_online,
                    )
                with action_cols[1]:
                    enhance_clicked = st.button(
                        "Enhance Prompt",
                        key="btn_enhance_prompt",
                        disabled=not query.strip() or not model_online,
                    )

                st.caption("Tip: Enhance the prompt to clarify intent before submitting.")

                if enhance_clicked:
                    with st.spinner(f"Enhancing prompt with {st.session_state.selected_model}..."):
                        success, message = enhance_query_prompt(query, st.session_state.selected_model)
                    if success:
                        st.session_state.enhanced_prompt = message
                        st.session_state.enhance_prompt_error = ""
                    else:
                        st.session_state.enhance_prompt_error = message
                        st.session_state.enhanced_prompt = ""

                if get_answer_clicked:
                    st.session_state._auto_submit = True

                if st.session_state.enhance_prompt_error:
                    st.error(st.session_state.enhance_prompt_error)

                if st.session_state.enhanced_prompt:
                    st.markdown("**Enhanced Prompt (edit before submitting)**")
                    enhanced_value = st.text_area(
                        "Enhanced Prompt",
                        value=st.session_state.enhanced_prompt,
                        key="enhanced_prompt_editor",
                        height=140,
                    )
                    enhanced_value = enhanced_value or ""
                    st.session_state.enhanced_prompt = enhanced_value

                    helper_cols = st.columns([1, 1])
                    with helper_cols[0]:
                        if st.button(
                            "Submit Enhanced Prompt",
                            type="primary",
                            key="btn_submit_enhanced",
                            disabled=not enhanced_value.strip() or not model_online,
                        ):
                            trimmed_prompt = enhanced_value.strip()
                            st.session_state.enhanced_prompt = ""
                            st.session_state.enhance_prompt_error = ""
                            st.session_state._autofill_query = trimmed_prompt
                            st.session_state._auto_submit = True
                            st.rerun()
                    with helper_cols[1]:
                        if st.button("Clear Enhanced Prompt", key="btn_clear_enhanced"):
                            st.session_state.enhanced_prompt = ""
                            st.session_state.enhance_prompt_error = ""
                            st.rerun()

                # If auto_submit flag is set (from suggested/followup click or Get Answer), execute the query
                if st.session_state.get('_auto_submit', False):
                    # Only proceed if there's a non-empty query and the model is online
                    if query.strip() and model_online:
                        # Clear the flag immediately so it doesn't loop
                        st.session_state._auto_submit = False
                        with st.spinner(f"Searching and generating answer using {st.session_state.selected_model}..."):
                            try:
                                selected_model = st.session_state.selected_model or "DeepSeek R1"

                                # Determine embedding model choice
                                use_legal_model = None  # Auto-detect by default
                                if "InLegalBERT" in st.session_state.embedding_model_choice:
                                    use_legal_model = True
                                elif "Sentence Transformers" in st.session_state.embedding_model_choice:
                                    use_legal_model = False

                                # Append the user message to chat history
                                st.session_state.chat_history.append({"role": "user", "content": query})

                                # Pass the recent chat history to the retriever so it can answer follow-ups
                                response = query_rag_system(
                                    query,
                                    index,
                                    embeddings,
                                    chunks,
                                    embedding_model,
                                    selected_model,
                                    use_legal_model=use_legal_model,
                                    conversation_history=st.session_state.chat_history,
                                    history_window=5,
                                )

                                # Response is expected to be a dict with 'answer' and 'followups'
                                answer_text = response['answer'] if isinstance(response, dict) else str(response)
                                followups = response.get('followups', []) if isinstance(response, dict) else []

                                # Append assistant message to chat history (include followups so UI can render buttons)
                                st.session_state.chat_history.append({"role": "assistant", "content": answer_text, "followups": followups})

                                # Also keep query_history (summary list) for admin view if present
                                st.session_state.query_history.append({
                                    'question': query,
                                    'answer': answer_text,
                                    'model': st.session_state.selected_model
                                })

                                st.session_state.enhanced_prompt = ""
                                st.session_state.enhance_prompt_error = ""

                                # Display answer in response section (chat rendering below)
                                st.markdown('<div class="response-section">', unsafe_allow_html=True)
                                st.subheader(f"Answer from {st.session_state.selected_model}:")
                                st.write(answer_text)

                                # Render follow-up questions as clickable buttons under this assistant message
                                if followups:
                                    st.write('---')
                                    st.write('You might also ask:')
                                    cols = st.columns(len(followups)) if len(followups) <= 4 else st.columns(4)
                                    for i, fq in enumerate(followups):
                                        col = cols[i % len(cols)]
                                        if col.button(fq, key=f"followup_{len(st.session_state.query_history)}_{i}"):
                                            # Auto-fill and auto-submit this follow-up question
                                            st.session_state._autofill_query = fq
                                            st.session_state._auto_submit = True
                                            st.rerun()

                                st.markdown('</div>', unsafe_allow_html=True)
                            except Exception as e:
                                st.error(f"Error generating answer: {e}")
                    elif not model_online:
                        st.warning("Please select an online model to ask questions.")
                    else:
                        st.warning("Please enter a question")
            else:
                st.error("Failed to load RAG components. Please refresh the knowledge base.")
        else:
            st.info("🔄 Upload documents to get started!")
            st.markdown("""
            **Getting Started:**
            1. Upload your documents (PDF, Excel, CSV, or Images)
            2. Click 'Process Files' to build the knowledge base
            3. Ask questions about your documents
            """)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # Chat history UI (render conversation with message alignment)
    if st.session_state.chat_history:
        st.divider()
        st.subheader("� Conversation")
        chat_container = st.container()

        with chat_container:
            for idx, msg in enumerate(st.session_state.chat_history):
                role = msg.get('role')
                content = msg.get('content', '')
                followups = msg.get('followups', []) if role == 'assistant' else []

                # Simple two-column layout to align messages left/right
                if role == 'user':
                    cols = st.columns([1, 4])
                    with cols[1]:
                        st.markdown(f"<div style='text-align: right; background:#DCF8C6; padding:10px; border-radius:10px; margin:6px; word-wrap:break-word;'>{content}</div>", unsafe_allow_html=True)
                else:
                    cols = st.columns([4, 1])
                    with cols[0]:
                        st.markdown(f"<div style='text-align: left; background:#F1F0F0; padding:10px; border-radius:10px; margin:6px; word-wrap:break-word;'>{content}</div>", unsafe_allow_html=True)

                        # Render followup buttons under assistant message
                        if followups:
                            btn_cols = st.columns(len(followups)) if len(followups) <= 4 else st.columns(4)
                            for j, fq in enumerate(followups):
                                with btn_cols[j % len(btn_cols)]:
                                    if st.button(fq, key=f"chat_followup_{idx}_{j}"):
                                        st.session_state._autofill_query = fq
                                        st.session_state._auto_submit = True
                                        st.rerun()
        

if __name__ == "__main__":
    main()
